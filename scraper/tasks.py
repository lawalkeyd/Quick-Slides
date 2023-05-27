from collections import defaultdict
import json

from bs4 import BeautifulSoup
from urllib.request import urlopen, Request
from pptx import Presentation
from pathlib import Path

from services.openai.gpt3 import summarize_text

class Scrape_Site:
    def __init__(self, site_url):
        self.filename = 'slide.log'
        self.site_url = site_url
        self.log = open(Path(__file__).parent / self.filename, 'w')
        self.log.write(self.site_url + '\n')
        self.log.write('Scraping Site ..' + '\n')
        req = Request(self.site_url, headers={'User-Agent': 'Mozilla/5.0'})
        html = urlopen(req).read()
        bs = BeautifulSoup(html, 'html.parser')
        self.booktitle = bs.find('title').text
        self.create_powerpoint()
        self.find_info(bs)
        self.slide.save(Path(__file__).parent.parent / 'slides_folder' / f'{self.booktitle}.pptx')
        self.log.close()

    def create_powerpoint(self):
        present_dir = Path(__file__).parent.parent / 'slides_folder'
        self.slide = Presentation(present_dir / 'Ion.pptx')

    def add_section(self, text):
        self.log.write('Starting section adding\n')
        slide_layout = self.slide.slide_layouts[0]
        slide = self.slide.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = text

    def add_slide(self, title, text):
        self.log.write('Starting slide adding\n')
        slide_layout = self.slide.slide_layouts[1]
        slide = self.slide.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = text

    def find_info(self, bs):
        def get_content(node):
            content = defaultdict(list)
            next_node = node.find_next_sibling()
            while next_node and next_node.name not in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                next_text = ''
                if next_node.name in ["p", "img"]:
                    if len(next_node.text.strip().split()) > 200:  # summarize if more than 200 words
                        next_text = summarize_text(next_node.text.strip())
                    content[node.name].append(next_text.strip() if next_node.name == "p" else next_node['src'])
                next_node = next_node.find_next_sibling()

            return dict(content)

        data = defaultdict(list)
        for heading in bs.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
            data[heading.name].append({
                "text": heading.text.strip(),
                "content": get_content(heading)
            })

        return json.dumps(data, indent=4)

       

Scrape_Site("https://en.wikipedia.org/wiki/Eddie_Murphy")